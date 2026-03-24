#!/usr/bin/env python3
"""Generate internal demo presentation PPTX — Hiring System + TOKAMAK Rally (English)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Colors
TOKAMAK_BLUE = RGBColor(0x2A, 0x72, 0xE5)
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xC8)
ACCENT_GREEN = RGBColor(0x00, 0xD1, 0x7A)
ACCENT_ORANGE = RGBColor(0xFF, 0x9F, 0x43)
ACCENT_RED = RGBColor(0xE6, 0x39, 0x46)
CARD_BG = RGBColor(0x1A, 0x25, 0x3C)
FONT = "Arial"
LOGO_PATH = "/Users/junwoong/.openclaw/workspace/tokamak-logo-ppt.jpg"

def add_bg(slide):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = DARK_BG

def add_logo(slide):
    pass  # Logo removed

def add_text(slide, text, left, top, width, height, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.color.rgb = color; p.font.bold = bold; p.font.name = FONT
    p.alignment = alignment
    return txBox

def add_bullets(slide, items, left, top, width, height, font_size=15, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(font_size); p.font.color.rgb = color; p.font.name = FONT
        p.space_after = Pt(6)
    return txBox

def card(slide, left, top, w, h, title, lines, accent=TOKAMAK_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = accent; shape.line.width = Pt(1.5)
    # accent bar removed
    add_text(slide, title, left+0.2, top+0.15, w-0.4, 0.4, font_size=18, color=accent, bold=True)
    add_bullets(slide, lines, left+0.2, top+0.6, w-0.4, h-0.8, font_size=13, color=LIGHT_GRAY)

def divider(slide, x, y, w):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = TOKAMAK_BLUE; line.line.fill.background()

# ═══════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, "TOKAMAK NETWORK", 1, 1.8, 14, 0.6, font_size=22, color=TOKAMAK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, "Internal Tools & Projects Demo", 1, 2.8, 14, 1, font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, "Developer Sourcing  ·  TOKAMAK Rally", 1, 5.0, 14, 0.6, font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
divider(slide, 5.5, 5.8, 5)
add_text(slide, "Presented by Jaden  |  March 2026", 1, 6.2, 14, 0.5, font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 2: Agenda
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_logo(slide)
add_text(slide, "AGENDA", 1, 0.5, 14, 0.6, font_size=36, color=WHITE, bold=True)
divider(slide, 1, 1.2, 2)

items = [
    ("01", "Developer Sourcing System", "GitHub-based talent discovery + LinkedIn outreach pipeline", ACCENT_GREEN),
    ("02", "TOKAMAK Rally", "Open-source blockchain racing game prototype", ACCENT_ORANGE),
    ("03", "Next Steps", "Roadmap and discussion", TOKAMAK_BLUE),
]
for i, (num, title, desc, accent) in enumerate(items):
    y = 2.2 + i * 1.8
    add_text(slide, num, 1.5, y, 1, 0.5, font_size=36, color=accent, bold=True)
    add_text(slide, title, 2.8, y, 8, 0.5, font_size=28, color=WHITE, bold=True)
    add_text(slide, desc, 2.8, y+0.5, 10, 0.4, font_size=18, color=LIGHT_GRAY)

# ═══════════════════════════════════════
# SLIDE 3: Developer Sourcing — How It Works
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_logo(slide)
add_text(slide, "01  DEVELOPER SOURCING SYSTEM", 1, 0.5, 12, 0.6, font_size=32, color=ACCENT_GREEN, bold=True)
add_text(slide, "How It Works", 1, 1.1, 12, 0.4, font_size=18, color=LIGHT_GRAY)

card(slide, 1, 1.8, 4.3, 3.3, "🔍 GitHub Scan Engine", [
    "Scans GitHub profiles using API",
    "Checks repos, languages, contributions",
    "Finds LinkedIn via Social Accounts API",
    "192 candidates found (61 LinkedIn + 131 GitHub)",
    "Auto-scoring with breakdown",
], ACCENT_GREEN)

card(slide, 5.8, 1.8, 4.3, 3.3, "🎯 Scoring & Matching", [
    "Domain alignment: 35% (highest weight)",
    "Execution ability: 25%",
    "Contactability: 15%",
    "Matches candidates against 14 team profiles",
    "Score breakdown on hover",
], TOKAMAK_BLUE)

card(slide, 10.6, 1.8, 4.3, 3.3, "📨 LinkedIn Outreach", [
    "Free connection requests (~100/week)",
    "For non-connections: paid plan required",
    "30 InMail messages per month",
    "Send 8-10 outreach messages per week",
    "15-day timeout → auto-reject",
], ACCENT_ORANGE)

# Current Status section
card(slide, 1, 5.6, 14, 2.5, "📊 Current Status", [
    "9 candidates contacted so far  →  4 declined  |  5 awaiting response (15-day timeout applies)",
    "No major results yet — this is expected. We continue this effort.",
    "This is NOT about filling open positions. We are building a long-term pipeline.",
    "In a fast-changing industry, we want to keep doors open to meet promising talent — even if the chance is small.",
], ACCENT_GREEN)

# ═══════════════════════════════════════
# SLIDE 4: TOKAMAK Rally Overview
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_logo(slide)
add_text(slide, "02  TOKAMAK RALLY", 1, 0.5, 12, 0.6, font_size=32, color=ACCENT_ORANGE, bold=True)
add_text(slide, "Open-Source Blockchain Racing Game", 1, 1.1, 12, 0.4, font_size=18, color=LIGHT_GRAY)

card(slide, 1, 1.7, 4.3, 3.2, "🎮 Game Overview", [
    "Phaser.js + Vite (runs in web browser)",
    "5 zones: Sand → Canyon → Riverbed → Mountain → Sprint",
    "5 cars including TOKAMAK GT",
    "121 waypoints / time-attack / drift mechanics",
    "Procedural audio — no sound files needed",
], ACCENT_ORANGE)

card(slide, 5.8, 1.7, 4.3, 3.2, "🔗 Blockchain Integration", [
    "MetaMask wallet connection",
    "Sepolia testnet leaderboard",
    "Guest play supported (no wallet needed)",
    "Tournament system with TOKAMAK token",
    "90% prize pool + 10% operator",
], TOKAMAK_BLUE)

card(slide, 10.6, 1.7, 4.3, 3.2, "📐 Design Principles", [
    "Open-source: Tokamak publishes, does not operate",
    "Deterministic obstacles (seeded RNG)",
    "Single linear track (no loops or forks)",
    "Hard difficulty — average users will struggle",
    "All original assets (no copyrighted material)",
], ACCENT_GREEN)

card(slide, 1, 5.4, 14, 2.8, "🎨 Graphics & Sound Highlights", [
    "258-color pixel art  |  20×28 car sprites with headlights/taillights  |  Zone-specific obstacles & particle effects",
    "Sprint zone: Tokamak Network banner walls on both sides + crowd behind  |  Finish line: checkered flag + spectators",
    "Procedural audio via Web Audio API: engine, drift screech (1350Hz band), offroad friction, Eurobeat BGM at 170 BPM",
    "Dust trail particles on unpaved roads  |  Real drift mechanic with body/movement direction split",
], ACCENT_ORANGE)

# ═══════════════════════════════════════
# SLIDE 5: TOKAMAK Rally — Future Plans
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_logo(slide)
add_text(slide, "TOKAMAK RALLY — FUTURE PLANS", 1, 0.5, 12, 0.6, font_size=28, color=ACCENT_ORANGE, bold=True)

card(slide, 1, 1.5, 7, 3.0, "🚀 Short-Term", [
    "Deploy on Vercel for internal team testing (~14 members)",
    "Collect feedback and adjust game balance",
    "Deploy RallyLeaderboard contract on Sepolia",
    "Iterate on difficulty and time budget",
], ACCENT_ORANGE)

card(slide, 8.5, 1.5, 6.5, 3.0, "🗺️ Feature Roadmap", [
    "New courses with different themes and layouts",
    "Map editor — custom track creation tool",
    "More car roster with unique characteristics",
    "Multiplayer ghost race mode",
], TOKAMAK_BLUE)

card(slide, 1, 5.0, 14, 3.0, "🔮 Long-Term Vision", [
    "Migrate from Sepolia testnet to Titan L2",
    "Full tournament system with TOKAMAK token integration (entry fee → prize pool)",
    "Community-created maps shared on-chain",
    "Open-source model: anyone can fork, run, and extend the game",
], ACCENT_GREEN)

# ═══════════════════════════════════════
# SLIDE 6: Next Steps (combined)
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_logo(slide)
add_text(slide, "03  NEXT STEPS", 1, 0.5, 12, 0.6, font_size=32, color=TOKAMAK_BLUE, bold=True)

card(slide, 1, 1.5, 7, 3.0, "🔍 Developer Sourcing", [
    "Continue weekly outreach (8-10 per week)",
    "Track response rates and refine messaging",
    "Purchase LinkedIn Sales Navigator Core",
    "Expand search criteria as needed",
], ACCENT_GREEN)

card(slide, 8.5, 1.5, 6.5, 3.0, "🎮 TOKAMAK Rally", [
    "Complete Vercel deployment",
    "Internal team playtesting session",
    "New course design + map editor prototype",
    "Sepolia smart contract deployment",
], ACCENT_ORANGE)

card(slide, 1, 5.0, 14, 2.5, "💡 Key Takeaway", [
    "Developer Sourcing: A low-cost, long-term effort to stay connected with potential talent in a changing landscape.",
    "TOKAMAK Rally: An open-source game that puts Tokamak Network's name in a fun, interactive format — ready for team testing soon.",
], TOKAMAK_BLUE)

# ═══════════════════════════════════════
# SLIDE 7: Q&A
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_logo(slide)
add_text(slide, "Q & A", 3, 2.5, 10, 1.2, font_size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
divider(slide, 6, 4.0, 4)
add_text(slide, "Thank You", 3, 4.5, 10, 0.6, font_size=28, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, "Tokamak Network  |  L2 On-Demand, Tailored Ethereum", 3, 5.3, 10, 0.5, font_size=16, color=TOKAMAK_BLUE, alignment=PP_ALIGN.CENTER)

# Save
out = "/Users/junwoong/.openclaw/workspace/tokamak-demo-presentation.pptx"
prs.save(out)
print(f"✅ Saved: {out}")
